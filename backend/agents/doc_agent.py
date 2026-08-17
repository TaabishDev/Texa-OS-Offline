import os
from typing import List, Dict, Any, Optional
from datetime import datetime

# Import document automation libraries
try:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
except ImportError:
    Document = None

try:
    from pptx import Presentation
    from pptx.util import Inches as PtInches, Pt as PtFont
except ImportError:
    Presentation = None

try:
    import openpyxl
    from openpyxl.styles import Font as xlFont, PatternFill, Alignment, Border, Side
except ImportError:
    openpyxl = None

try:
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
except ImportError:
    SimpleDocTemplate = None

class DocAgent:
    def __init__(self):
        self.output_dir = os.path.expanduser("~/Downloads")
        os.makedirs(self.output_dir, exist_ok=True)

    def generate(self, doc_type: str, title: str, content: List[Dict[str, Any]], file_name: Optional[str] = None) -> str:
        doc_type = doc_type.lower()
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        if not file_name:
            safe_title = "".join([c if c.isalnum() else "_" for c in title]).strip("_")
            file_name = f"{safe_title}_{timestamp}"
            
        if doc_type == "word" or doc_type == "docx":
            if not file_name.endswith(".docx"):
                file_name += ".docx"
            filepath = os.path.join(self.output_dir, file_name)
            self._generate_docx(title, content, filepath)
            
        elif doc_type == "ppt" or doc_type == "pptx":
            if not file_name.endswith(".pptx"):
                file_name += ".pptx"
            filepath = os.path.join(self.output_dir, file_name)
            self._generate_pptx(title, content, filepath)
            
        elif doc_type == "excel" or doc_type == "xlsx":
            if not file_name.endswith(".xlsx"):
                file_name += ".xlsx"
            filepath = os.path.join(self.output_dir, file_name)
            self._generate_xlsx(title, content, filepath)
            
        elif doc_type == "pdf":
            if not file_name.endswith(".pdf"):
                file_name += ".pdf"
            filepath = os.path.join(self.output_dir, file_name)
            self._generate_pdf(title, content, filepath)
            
        else:
            raise ValueError(f"Unsupported document format: {doc_type}")
            
        return filepath

    def _generate_docx(self, title: str, content: List[Dict[str, Any]], filepath: str):
        if not Document:
            raise RuntimeError("python-docx is not installed. Install requirements first.")
            
        doc = Document()
        
        # Styles setup (Premium layout: Teal & Charcoal)
        styles = doc.styles
        
        # Main Title (Cover design)
        t_para = doc.add_paragraph()
        t_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        t_run = t_para.add_run(title.upper())
        t_run.font.name = "Arial"
        t_run.font.size = Pt(26)
        t_run.font.bold = True
        t_run.font.color.rgb = RGBColor(11, 18, 48)  # Deep Indigo
        
        doc.add_paragraph().add_run("Created by TEXA - AI Executive Assistant").italic = True
        doc.add_paragraph(datetime.now().strftime("%B %d, %Y"))
        doc.add_page_break()
        
        for block in content:
            block_type = block.get("type", "paragraph")
            value = block.get("value", "")
            
            if block_type == "heading_1":
                h = doc.add_heading(level=1)
                run = h.add_run(value)
                run.font.color.rgb = RGBColor(11, 18, 48)
                run.font.size = Pt(20)
                
            elif block_type == "heading_2":
                h = doc.add_heading(level=2)
                run = h.add_run(value)
                run.font.color.rgb = RGBColor(46, 117, 182)
                run.font.size = Pt(15)
                
            elif block_type == "paragraph":
                p = doc.add_paragraph(value)
                p.paragraph_format.line_spacing = 1.15
                
            elif block_type == "bullet":
                doc.add_paragraph(value, style='List Bullet')
                
            elif block_type == "table":
                # Expecting value as a list of lists: [[col1, col2], [val1, val2]]
                rows_data = value
                if not isinstance(rows_data, list) or not rows_data:
                    continue
                num_cols = len(rows_data[0])
                num_rows = len(rows_data)
                
                table = doc.add_table(rows=num_rows, cols=num_cols)
                table.style = 'Light Shading Accent 1'
                
                for r_idx, row_data in enumerate(rows_data):
                    row = table.rows[r_idx]
                    for c_idx, cell_val in enumerate(row_data):
                        cell = row.cells[c_idx]
                        cell.text = str(cell_val)
            
            doc.add_paragraph("")  # spacing
            
        doc.save(filepath)

    def _generate_pptx(self, title: str, content: List[Dict[str, Any]], filepath: str):
        if not Presentation:
            raise RuntimeError("python-pptx is not installed. Install requirements first.")
            
        prs = Presentation()
        
        # 1. Slide Title Layout (0)
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Generated by Texa\nAbdul's AI Operating System"
        
        # 2. Add Content Slides
        current_slide = None
        current_tf = None
        
        for block in content:
            block_type = block.get("type", "paragraph")
            value = block.get("value", "")
            
            # Start new slide on Heading 1
            if block_type == "heading_1":
                current_slide = prs.slides.add_slide(prs.slide_layouts[1]) # Title & Content
                current_slide.shapes.title.text = value
                # Get the body textbox placeholder
                body_shape = current_slide.placeholders[1]
                current_tf = body_shape.text_frame
                # Clear default paragraph
                current_tf.text = ""
                
            elif block_type == "bullet" or block_type == "paragraph":
                if not current_slide:
                    # Create a default slide if heading_1 wasn't specified yet
                    current_slide = prs.slides.add_slide(prs.slide_layouts[1])
                    current_slide.shapes.title.text = "Overview"
                    body_shape = current_slide.placeholders[1]
                    current_tf = body_shape.text_frame
                    current_tf.text = ""
                    
                p = current_tf.add_paragraph()
                p.text = value
                p.font.size = PtFont(16)
                if block_type == "bullet":
                    p.level = 0
                    
            elif block_type == "table":
                # Skip tables in simple PPT generation for now, or print as string
                if current_tf:
                    p = current_tf.add_paragraph()
                    p.text = "[Table content generated - view main report for complete data]"
                    p.font.italic = True
                    
        prs.save(filepath)

    def _generate_xlsx(self, title: str, content: List[Dict[str, Any]], filepath: str):
        if not openpyxl:
            raise RuntimeError("openpyxl is not installed. Install requirements first.")
            
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "TEXA Report"
        
        # Style Definitions
        title_font = xlFont(name="Arial", size=16, bold=True, color="000B1230")
        header_font = xlFont(name="Arial", size=11, bold=True, color="FFFFFFFF")
        header_fill = PatternFill(start_color="FF0B1230", end_color="FF0B1230", fill_type="solid")
        border_side = Side(border_style="thin", color="FFD9D9D9")
        thin_border = Border(left=border_side, right=border_side, top=border_side, bottom=border_side)
        
        # Add Title Row
        ws.merge_cells("A1:E1")
        ws["A1"] = title
        ws["A1"].font = title_font
        ws.row_dimensions[1].height = 35
        ws["A1"].alignment = Alignment(vertical="center")
        
        row_cursor = 3
        
        for block in content:
            block_type = block.get("type", "paragraph")
            value = block.get("value", "")
            
            if block_type == "heading_1" or block_type == "heading_2":
                ws.cell(row=row_cursor, column=1, value=value).font = xlFont(name="Arial", size=12, bold=True, color="002E75B6")
                row_cursor += 2
                
            elif block_type == "table":
                # Expect list of lists
                rows_data = value
                if not isinstance(rows_data, list) or not rows_data:
                    continue
                    
                for r_idx, row_data in enumerate(rows_data):
                    for c_idx, cell_val in enumerate(row_data):
                        cell = ws.cell(row=row_cursor + r_idx, column=c_idx + 1, value=cell_val)
                        cell.border = thin_border
                        
                        # Apply header styles to first row of table
                        if r_idx == 0:
                            cell.font = header_font
                            cell.fill = header_fill
                            cell.alignment = Alignment(horizontal="center")
                            
                    ws.row_dimensions[row_cursor + r_idx].height = 20
                    
                row_cursor += len(rows_data) + 2
                
            elif block_type == "paragraph" or block_type == "bullet":
                ws.cell(row=row_cursor, column=1, value=value).font = xlFont(name="Arial", size=10)
                row_cursor += 1
                
        # Auto-adjust column widths
        for col in ws.columns:
            max_len = 0
            for cell in col:
                if cell.value:
                    max_len = max(max_len, len(str(cell.value)))
            col_letter = openpyxl.utils.get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = max(max_len + 3, 12)
            
        wb.save(filepath)

    def _generate_pdf(self, title: str, content: List[Dict[str, Any]], filepath: str):
        if not SimpleDocTemplate:
            raise RuntimeError("reportlab is not installed. Install requirements first.")
            
        doc = SimpleDocTemplate(filepath, pagesize=letter, rightMargin=54, leftMargin=54, topMargin=54, bottomMargin=54)
        story = []
        
        # Styles
        styles = getSampleStyleSheet()
        
        title_style = ParagraphStyle(
            'ReportTitle',
            parent=styles['Heading1'],
            fontName='Helvetica-Bold',
            fontSize=24,
            textColor=colors.HexColor('#0b1230'),
            spaceAfter=30,
            alignment=1 # Center
        )
        
        h1_style = ParagraphStyle(
            'ReportH1',
            parent=styles['Heading2'],
            fontName='Helvetica-Bold',
            fontSize=16,
            textColor=colors.HexColor('#0b1230'),
            spaceBefore=15,
            spaceAfter=10
        )
        
        h2_style = ParagraphStyle(
            'ReportH2',
            parent=styles['Heading3'],
            fontName='Helvetica-Bold',
            fontSize=12,
            textColor=colors.HexColor('#2e75b6'),
            spaceBefore=10,
            spaceAfter=8
        )
        
        body_style = ParagraphStyle(
            'ReportBody',
            parent=styles['Normal'],
            fontName='Helvetica',
            fontSize=10.5,
            leading=14,
            textColor=colors.HexColor('#222222'),
            spaceAfter=10
        )
        
        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 15))
        story.append(Paragraph(f"Date: {datetime.now().strftime('%B %d, %Y')}", body_style))
        story.append(Paragraph("System Engine: TEXA Assistant Platform", body_style))
        story.append(Spacer(1, 20))
        
        for block in content:
            block_type = block.get("type", "paragraph")
            value = block.get("value", "")
            
            if block_type == "heading_1":
                story.append(Paragraph(value, h1_style))
            elif block_type == "heading_2":
                story.append(Paragraph(value, h2_style))
            elif block_type == "paragraph":
                story.append(Paragraph(value, body_style))
            elif block_type == "bullet":
                story.append(Paragraph(f"• {value}", body_style))
            elif block_type == "table":
                rows_data = value
                if not isinstance(rows_data, list) or not rows_data:
                    continue
                
                # Make cells Paragraphs so they wrap nicely
                formatted_table_data = []
                for r_idx, row in enumerate(rows_data):
                    formatted_row = []
                    for col in row:
                        if r_idx == 0:
                            # Table Header
                            th_style = ParagraphStyle('TH', parent=body_style, fontName='Helvetica-Bold', textColor=colors.white)
                            formatted_row.append(Paragraph(str(col), th_style))
                        else:
                            formatted_row.append(Paragraph(str(col), body_style))
                    formatted_table_data.append(formatted_row)
                
                t = Table(formatted_table_data)
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#0b1230')),
                    ('ALIGN', (0,0), (-1,-1), 'LEFT'),
                    ('VALIGN', (0,0), (-1,-1), 'TOP'),
                    ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cccccc')),
                    ('TOPPADDING', (0,0), (-1,-1), 6),
                    ('BOTTOMPADDING', (0,0), (-1,-1), 6),
                ]))
                story.append(Spacer(1, 5))
                story.append(t)
                story.append(Spacer(1, 10))
                
        doc.build(story)
